"""Generate curriculum slideshow presentation for The AI Collaborator."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- color palette ----------
BG_DARK      = RGBColor(0x16, 0x1b, 0x22)
BG_CARD      = RGBColor(0x1e, 0x24, 0x2e)
TEXT_PRIMARY  = RGBColor(0xE6, 0xED, 0xF3)
TEXT_SECONDARY= RGBColor(0x8B, 0x94, 0x9E)
ACCENT_GREEN  = RGBColor(0x4A, 0x90, 0x79)
ACCENT_PURPLE = RGBColor(0x90, 0x79, 0xB0)
ACCENT_BROWN  = RGBColor(0xB0, 0x80, 0x50)
ACCENT_BLUE   = RGBColor(0x50, 0x90, 0xB0)
ACCENT_ORANGE = RGBColor(0xB0, 0x70, 0x50)
ACCENT_SLATE  = RGBColor(0x60, 0x70, 0x90)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT     = RGBColor(0x58, 0xA6, 0xFF)

MODULE_COLORS = {
    'Foundation':            RGBColor(0x4A, 0x90, 0x79),
    'Documentation & Trust': RGBColor(0x90, 0x79, 0xB0),
    'Workflow':              RGBColor(0xB0, 0x80, 0x50),
    'Advanced':              RGBColor(0x50, 0x90, 0xB0),
}

CONCEPT_COLORS = {
    'Context Assembly':     ACCENT_GREEN,
    'Quality Judgment':     ACCENT_PURPLE,
    'Task Decomposition':   ACCENT_BROWN,
    'Iterative Refinement': ACCENT_BLUE,
    'Workflow Integration': ACCENT_ORANGE,
    'Frontier Recognition': ACCENT_SLATE,
}

# ---------- data ----------
CONCEPTS = [
    {
        'name': 'Context Assembly',
        'tagline': 'Curating the briefing that shapes AI output quality',
        'description': 'Gathering and presenting relevant background information to enable effective AI collaboration. Identifying what information is needed, organizing it coherently, and providing sufficient detail without overwhelming noise.',
        'lessons': [1, 3, 4],
    },
    {
        'name': 'Quality Judgment',
        'tagline': 'Distinguishing "looks right" from "is right"',
        'description': 'Critically evaluating AI-generated outputs for accuracy, completeness, and fitness for purpose. Recognizing errors, identifying gaps, and determining whether the output genuinely meets the intended need.',
        'lessons': [5, 6],
    },
    {
        'name': 'Task Decomposition',
        'tagline': 'Breaking complex problems into AI-appropriate chunks',
        'description': 'Breaking complex problems into smaller components that can be addressed sequentially or in parallel. Understanding which subtasks are AI-appropriate and where human judgment is required.',
        'lessons': [7, 8],
    },
    {
        'name': 'Iterative Refinement',
        'tagline': 'Steering toward outcomes through successive approximations',
        'description': 'Progressively improving outputs through cycles of feedback, adjustment, and revision. Knowing what to ask for, how to redirect, and when "good enough" has been reached.',
        'lessons': [2, 9],
    },
    {
        'name': 'Workflow Integration',
        'tagline': 'Embedding AI into sustainable work patterns',
        'description': 'Embedding AI collaboration into existing work processes in sustainable, practical ways. Identifying where AI adds genuine value and creating repeatable patterns.',
        'lessons': [10, 12],
    },
    {
        'name': 'Frontier Recognition',
        'tagline': 'Knowing the boundaries of AI capability',
        'description': 'Understanding the current boundaries of AI capability—what it can and cannot do reliably, where it excels versus struggles, and how those boundaries are shifting.',
        'lessons': [11],
    },
]

MODULES = [
    {
        'name': 'Foundation',
        'lessons': [
            {'num': 1, 'title': 'Context Tracker', 'concept': 'Context Assembly',
             'problem': 'AI conversations fail when critical context is missing. You waste time on back-and-forth clarifications or get unusable outputs because you forgot to mention key constraints.',
             'skill': 'Identify your personal context gaps by analyzing past conversations. Discover what information you consistently forget to provide so you can fix it upfront.'},
            {'num': 2, 'title': 'Feedback Analyzer', 'concept': 'Iterative Refinement',
             'problem': 'Vague feedback like "make it better" wastes iteration cycles. Without specific, actionable feedback, you keep going in circles.',
             'skill': 'Write feedback that identifies specific locations, states clear actions, and explains reasoning. Learn to spot vague patterns in your own feedback.'},
            {'num': 3, 'title': 'Template Builder', 'concept': 'Context Assembly',
             'problem': 'You keep forgetting to provide the same context over and over. Each conversation starts from scratch, wasting time re-explaining your project and constraints.',
             'skill': 'Build reusable templates that capture the context AI needs upfront. Turn your insights into structured prompts you can use consistently.'},
        ],
    },
    {
        'name': 'Documentation & Trust',
        'lessons': [
            {'num': 4, 'title': 'Context Docs', 'concept': 'Context Assembly',
             'problem': 'Every new AI session starts from scratch. You waste time re-explaining project context, and the AI makes the same mistakes you\'ve already corrected.',
             'skill': 'Maintain living context documents that capture project state, decisions, issues, and lessons. Start each session with full context for immediate productivity.'},
            {'num': 5, 'title': 'Trust Matrix', 'concept': 'Quality Judgment',
             'problem': 'You either over-verify everything (wasting time) or blindly trust AI output (introducing errors). Without calibrated judgment, you can\'t efficiently allocate review effort.',
             'skill': 'Build a personal trust matrix by tracking predictions about AI accuracy. Learn which output types you can trust and which require careful verification.'},
            {'num': 6, 'title': 'Verification Tools', 'concept': 'Quality Judgment',
             'problem': 'Without systematic verification, you either waste time over-checking or miss critical errors by under-checking outputs that needed scrutiny.',
             'skill': 'Create reusable verification checklists tied to output types. Track which checks actually catch issues to refine your process over time.'},
        ],
    },
    {
        'name': 'Workflow',
        'lessons': [
            {'num': 7, 'title': 'Task Decomposer', 'concept': 'Task Decomposition',
             'problem': 'Without decomposition skills, you either delegate tasks that need your judgment (poor results) or do everything yourself (wasting AI\'s potential).',
             'skill': 'Break projects into subtasks and categorize each as AI-Optimal, Collaborative, or Human-Primary. Sequence tasks with dependencies to optimize the human-AI division of labor.'},
            {'num': 8, 'title': 'Delegation Tracker', 'concept': 'Task Decomposition',
             'problem': 'Knowing what to delegate is only half the battle. Without structured delegation, you give vague instructions and get disappointing results.',
             'skill': 'Create delegation templates with clear context, objectives, scope, deliverables, and success criteria. Execute decomposed tasks while tracking outcomes.'},
            {'num': 9, 'title': 'Iteration Passes', 'concept': 'Iterative Refinement',
             'problem': 'Random iteration ("make it better") wastes cycles and leads to scope creep. Without structure, you keep tweaking without knowing when "done" is reached.',
             'skill': 'Use the 70-85-95 framework to iterate with purpose. Each pass has a specific focus and key question, so you know exactly what to evaluate and when to move on.'},
            {'num': 10, 'title': 'Status Reporter', 'concept': 'Workflow Integration',
             'problem': 'Recurring tasks like status reports, meeting summaries, and client updates eat up valuable time when done manually each time.',
             'skill': 'Design AI-integrated workflows for recurring tasks. Create templates, track inputs, and measure time savings to build sustainable AI collaboration habits.'},
        ],
    },
    {
        'name': 'Advanced',
        'lessons': [
            {'num': 11, 'title': 'Frontier Mapper', 'concept': 'Frontier Recognition',
             'problem': 'AI capabilities change rapidly. Without tracking the frontier, you can\'t anticipate what\'s possible or prepare for new opportunities.',
             'skill': 'Map AI reliability zones and log frontier encounters to build your personal AI capability map. Track where AI excels, where it struggles, and where the boundaries are shifting.'},
            {'num': 12, 'title': 'Reference Card', 'concept': 'Workflow Integration',
             'problem': 'All your learnings across lessons can be overwhelming to remember and apply consistently in daily work.',
             'skill': 'Generate your personal AI collaboration quick reference card from your learnings across all lessons. Synthesize your insights into an actionable cheat sheet.'},
        ],
    },
]

COMPARISON = {
    'typical': [
        'Prompt tricks without judgment to apply them',
        'Tool training that expires when the UI changes',
        'One-shot demos that don\'t stick',
        'No framework for when AI fails',
    ],
    'platform': [
        'Transferable managerial skills',
        'Works with any AI tool',
        'Practice-based habit building',
        'Know exactly when to trust and verify',
    ],
}

# ---------- helpers ----------
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def in_to_emu(inches):
    return int(inches * 914400)

# ---------- slide builders ----------
def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Accent bar at top
    add_accent_bar(slide, 0, 0, in_to_emu(13.333), in_to_emu(0.06), ACCENT_GREEN)

    add_textbox(slide, in_to_emu(1.5), in_to_emu(2.2), in_to_emu(10.333), in_to_emu(1.2),
                'The AI Collaborator', font_size=44, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, in_to_emu(2), in_to_emu(3.5), in_to_emu(9.333), in_to_emu(0.8),
                'Learn to manage AI — not just use it', font_size=24,
                color=HIGHLIGHT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, in_to_emu(2.5), in_to_emu(4.6), in_to_emu(8.333), in_to_emu(1),
                '12 hands-on lessons  •  6 managerial skills  •  4 progressive modules',
                font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


def build_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, in_to_emu(13.333), in_to_emu(0.06), ACCENT_GREEN)

    add_textbox(slide, in_to_emu(1), in_to_emu(0.6), in_to_emu(11.333), in_to_emu(0.8),
                'Most AI training teaches the wrong thing', font_size=32,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Two columns
    col_w = in_to_emu(5)
    left_x = in_to_emu(1)
    right_x = in_to_emu(7.333)

    # Left column — Typical
    add_rounded_rect(slide, left_x, in_to_emu(1.8), col_w, in_to_emu(4.5), RGBColor(0x2a, 0x1a, 0x1a))
    add_textbox(slide, left_x + in_to_emu(0.3), in_to_emu(2.0), col_w - in_to_emu(0.6), in_to_emu(0.5),
                'Typical AI Training', font_size=20, color=RGBColor(0xFF, 0x80, 0x80), bold=True)
    for i, item in enumerate(COMPARISON['typical']):
        add_textbox(slide, left_x + in_to_emu(0.3), in_to_emu(2.7 + i * 0.7),
                    col_w - in_to_emu(0.6), in_to_emu(0.6),
                    f'✗  {item}', font_size=15, color=TEXT_SECONDARY)

    # Right column — This platform
    add_rounded_rect(slide, right_x, in_to_emu(1.8), col_w, in_to_emu(4.5), RGBColor(0x1a, 0x2a, 0x1a))
    add_textbox(slide, right_x + in_to_emu(0.3), in_to_emu(2.0), col_w - in_to_emu(0.6), in_to_emu(0.5),
                'This Platform', font_size=20, color=RGBColor(0x80, 0xFF, 0x80), bold=True)
    for i, item in enumerate(COMPARISON['platform']):
        add_textbox(slide, right_x + in_to_emu(0.3), in_to_emu(2.7 + i * 0.7),
                    col_w - in_to_emu(0.6), in_to_emu(0.6),
                    f'✓  {item}', font_size=15, color=TEXT_PRIMARY)


def build_concepts_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, in_to_emu(13.333), in_to_emu(0.06), ACCENT_GREEN)

    add_textbox(slide, in_to_emu(1), in_to_emu(0.5), in_to_emu(11.333), in_to_emu(0.7),
                'Six Skills That Make the Difference', font_size=32,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 3x2 grid of concept cards
    card_w = in_to_emu(3.5)
    card_h = in_to_emu(2.6)
    gap_x = in_to_emu(0.4)
    gap_y = in_to_emu(0.35)
    start_x = in_to_emu(0.85)
    start_y = in_to_emu(1.5)

    for idx, c in enumerate(CONCEPTS):
        col = idx % 3
        row = idx // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        accent = CONCEPT_COLORS[c['name']]
        add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD)
        add_accent_bar(slide, x, y, in_to_emu(0.06), card_h, accent)

        add_textbox(slide, x + in_to_emu(0.25), y + in_to_emu(0.2),
                    card_w - in_to_emu(0.4), in_to_emu(0.4),
                    c['name'], font_size=16, color=accent, bold=True)
        add_textbox(slide, x + in_to_emu(0.25), y + in_to_emu(0.6),
                    card_w - in_to_emu(0.4), in_to_emu(0.4),
                    c['tagline'], font_size=11, color=HIGHLIGHT)

        # Lessons tag
        lessons_str = ', '.join([f'L{n}' for n in c['lessons']])
        add_textbox(slide, x + in_to_emu(0.25), y + card_h - in_to_emu(0.5),
                    card_w - in_to_emu(0.4), in_to_emu(0.35),
                    f'Lessons: {lessons_str}', font_size=10, color=TEXT_SECONDARY)


def build_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, in_to_emu(13.333), in_to_emu(0.06), ACCENT_GREEN)

    add_textbox(slide, in_to_emu(1), in_to_emu(0.8), in_to_emu(11.333), in_to_emu(0.7),
                'How It Works', font_size=32, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    steps = [
        ('1', 'Learn', 'Each lesson introduces a concept with clear definitions and real-world context.', ACCENT_GREEN),
        ('2', 'Practice', 'Interactive exercises analyze your actual AI conversations and workflows.', ACCENT_BLUE),
        ('3', 'Track', 'Personal dashboards show your growth across all six skills over time.', ACCENT_PURPLE),
    ]

    card_w = in_to_emu(3.3)
    start_x = in_to_emu(1.2)
    gap = in_to_emu(0.6)

    for i, (num, title, desc, color) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        y = in_to_emu(2.2)

        add_rounded_rect(slide, x, y, card_w, in_to_emu(3.5), BG_CARD)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(x + card_w // 2 - in_to_emu(0.4)),
            Emu(y + in_to_emu(0.3)), Emu(in_to_emu(0.8)), Emu(in_to_emu(0.8))
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False

        add_textbox(slide, x + in_to_emu(0.3), y + in_to_emu(1.3),
                    card_w - in_to_emu(0.6), in_to_emu(0.5),
                    title, font_size=22, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + in_to_emu(0.3), y + in_to_emu(2.0),
                    card_w - in_to_emu(0.6), in_to_emu(1.2),
                    desc, font_size=14, color=TEXT_SECONDARY,
                    alignment=PP_ALIGN.CENTER)


def build_module_title(prs, module_name, lesson_count):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    color = MODULE_COLORS[module_name]
    add_accent_bar(slide, 0, 0, in_to_emu(13.333), in_to_emu(0.06), color)

    # Module label
    add_textbox(slide, in_to_emu(1), in_to_emu(2.5), in_to_emu(11.333), in_to_emu(0.5),
                'MODULE', font_size=14, color=TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER, bold=True)

    add_textbox(slide, in_to_emu(1), in_to_emu(3.0), in_to_emu(11.333), in_to_emu(1),
                module_name, font_size=40, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, in_to_emu(1), in_to_emu(4.2), in_to_emu(11.333), in_to_emu(0.5),
                f'{lesson_count} {"lesson" if lesson_count == 1 else "lessons"}',
                font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


def build_lesson_slide(prs, lesson, module_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    mod_color = MODULE_COLORS[module_name]
    concept_color = CONCEPT_COLORS[lesson['concept']]
    add_accent_bar(slide, 0, 0, in_to_emu(13.333), in_to_emu(0.06), mod_color)

    # Lesson number + title
    add_textbox(slide, in_to_emu(0.8), in_to_emu(0.4), in_to_emu(6), in_to_emu(0.4),
                f'LESSON {lesson["num"]}', font_size=13, color=TEXT_SECONDARY, bold=True)

    add_textbox(slide, in_to_emu(0.8), in_to_emu(0.8), in_to_emu(8), in_to_emu(0.7),
                lesson['title'], font_size=30, color=WHITE, bold=True)

    # Concept badge
    badge = add_rounded_rect(slide, in_to_emu(0.8), in_to_emu(1.6), in_to_emu(2.5), in_to_emu(0.4), concept_color)
    add_textbox(slide, in_to_emu(0.85), in_to_emu(1.6), in_to_emu(2.5), in_to_emu(0.4),
                lesson['concept'], font_size=11, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # The Problem card
    problem_y = in_to_emu(2.4)
    card_w = in_to_emu(11.733)
    add_rounded_rect(slide, in_to_emu(0.8), problem_y, card_w, in_to_emu(1.9), BG_CARD)
    add_accent_bar(slide, in_to_emu(0.8), problem_y, in_to_emu(0.06), in_to_emu(1.9), RGBColor(0xFF, 0x80, 0x80))
    add_textbox(slide, in_to_emu(1.1), problem_y + in_to_emu(0.15), in_to_emu(3), in_to_emu(0.4),
                'The Problem', font_size=14, color=RGBColor(0xFF, 0x80, 0x80), bold=True)
    add_textbox(slide, in_to_emu(1.1), problem_y + in_to_emu(0.55), card_w - in_to_emu(0.6), in_to_emu(1.2),
                lesson['problem'], font_size=14, color=TEXT_SECONDARY)

    # The Skill card
    skill_y = in_to_emu(4.6)
    add_rounded_rect(slide, in_to_emu(0.8), skill_y, card_w, in_to_emu(1.9), BG_CARD)
    add_accent_bar(slide, in_to_emu(0.8), skill_y, in_to_emu(0.06), in_to_emu(1.9), ACCENT_GREEN)
    add_textbox(slide, in_to_emu(1.1), skill_y + in_to_emu(0.15), in_to_emu(3), in_to_emu(0.4),
                'The Skill', font_size=14, color=ACCENT_GREEN, bold=True)
    add_textbox(slide, in_to_emu(1.1), skill_y + in_to_emu(0.55), card_w - in_to_emu(0.6), in_to_emu(1.2),
                lesson['skill'], font_size=14, color=TEXT_PRIMARY)


def build_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_bar(slide, 0, 0, in_to_emu(13.333), in_to_emu(0.06), ACCENT_GREEN)

    add_textbox(slide, in_to_emu(1), in_to_emu(1.5), in_to_emu(11.333), in_to_emu(1),
                'Build the skills that compound', font_size=36,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, in_to_emu(2), in_to_emu(2.8), in_to_emu(9.333), in_to_emu(1),
                'All 12 lessons available from day one.\nNo sequential unlocking.\nStart wherever your skills need the most work.',
                font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Stats row
    stats = [('12', 'Lessons'), ('6', 'Core Skills'), ('4', 'Modules')]
    start_x = in_to_emu(2.5)
    gap = in_to_emu(3.3)
    for i, (num, label) in enumerate(stats):
        x = start_x + i * gap
        y = in_to_emu(4.5)
        add_textbox(slide, x, y, in_to_emu(2), in_to_emu(0.7),
                    num, font_size=40, color=HIGHLIGHT, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + in_to_emu(0.7), in_to_emu(2), in_to_emu(0.4),
                    label, font_size=16, color=TEXT_SECONDARY,
                    alignment=PP_ALIGN.CENTER)


# ---------- main ----------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)
    build_problem_slide(prs)
    build_concepts_overview(prs)
    build_how_it_works(prs)

    for module in MODULES:
        build_module_title(prs, module['name'], len(module['lessons']))
        for lesson in module['lessons']:
            build_lesson_slide(prs, lesson, module['name'])

    build_summary_slide(prs)

    out_path = '/root/ClaudeProjects/ai-manager-skills/AI_Collaborator_Curriculum.pptx'
    prs.save(out_path)
    print(f'Saved: {out_path}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
